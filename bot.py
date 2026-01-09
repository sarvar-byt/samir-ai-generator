import streamlit as st
from pptx import Presentation
from docx import Document
import io
import telebot
from telebot import types
import threading

# ==========================================
# 1. SOZLAMALAR
# ==========================================
BOT_TOKEN = "8224322194:AAFyS95idEkFCJPT-hMMC_mx9jzNGiqxNW8"
MY_CHAT_ID = "5588598964" 

bot = telebot.TeleBot(BOT_TOKEN)

# ==========================================
# 2. UMUMIY GENERATSIYA FUNKSIYALARI
# ==========================================

def get_rich_content(topic):
    """Mavzu bo'yicha kengaytirilgan matn yaratish"""
    return [
        {"h": f"{topic}: Kirish", "c": f"{topic} tushunchasining kelib chiqishi va uning bugungi kundagi ahamiyati haqida batafsil ma'lumot. Ushbu soha jadal rivojlanmoqda."},
        {"h": "Asosiy maqsad va vazifalar", "c": f"{topic} yo'nalishidagi asosiy strategik maqsadlar va ularni amalga oshirish bosqichlari. Bu jarayon ko'plab resurslarni talab etadi."},
        {"h": "Muammolar va Tahlil", "c": f"Hozirgi vaqtda {topic} sohasida uchrayotgan to'siqlar: texnik yetishmovchiliklar, iqtisodiy omillar va inson resurslari tahlili."},
        {"h": "Innovatsion Yechimlar", "c": f"{topic}ni rivojlantirish uchun yangi texnologiyalar va innovatsion yondashuvlarni qo'llash imkoniyatlari o'rganilmoqda."},
        {"h": "Xulosa va Istiqbollar", "c": f"Xulosa qilib aytganda, {topic} kelajakda iqtisodiyot va jamiyatning ajralmas qismiga aylanadi. Kutilayotgan natijalar ijobiy."}
    ]

def create_pptx(topic):
    prs = Presentation()
    content = get_rich_content(topic)
    for data in content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = data['h']
        slide.placeholders[1].text = data['c']
    
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()

def create_docx(topic):
    doc = Document()
    doc.add_heading(topic.upper(), 0)
    content = get_rich_content(topic)
    for data in content:
        doc.add_heading(data['h'], level=1)
        doc.add_paragraph(data['c'])
        doc.add_paragraph("-" * 20)
    
    output = io.BytesIO()
    doc.save(output)
    return output.getvalue()

# ==========================================
# 3. TELEGRAM BOT LOGIKASI (ODAMLAR UCHUN)
# ==========================================

@bot.message_handler(commands=['start'])
def send_welcome(message):
    # Foydalanuvchiga start bosganda doimiy tugmani ham ko'rsatamiz
    markup = types.ReplyKeyboardMarkup(resize_keyboard=True)
    markup.add("🔄 Yangi mavzu yozish")
    bot.send_message(message.chat.id, "Salom! Men Samir AI botiman. 🤖\nPrezentatsiya yoki Word yasash uchun mavzu nomini yuboring:", reply_markup=markup)

@bot.message_handler(func=lambda message: True)
def handle_message(message):
    if message.text == "🔄 Yangi mavzu yozish":
        bot.send_message(message.chat.id, "Marhamat, yangi mavzu nomini kiriting:")
        return

    topic = message.text
    markup = types.InlineKeyboardMarkup()
    btn1 = types.InlineKeyboardButton("📊 PPTX (Prezentatsiya)", callback_data=f"pptx|{topic}")
    btn2 = types.InlineKeyboardButton("📝 DOCX (Word)", callback_data=f"docx|{topic}")
    markup.add(btn1, btn2)
    bot.send_message(message.chat.id, f"'{topic}' mavzusi bo'yicha qaysi formatda hujjat kerak?", reply_markup=markup)

@bot.callback_query_handler(func=lambda call: True)
def callback_inline(call):
    type_file, topic = call.data.split("|")
    bot.edit_message_text(chat_id=call.message.chat.id, message_id=call.message.message_id, text=f"⌛ '{topic}' bo'yicha {type_file} tayyorlanmoqda...")
    
    if type_file == "pptx":
        file_bytes = create_pptx(topic)
        bot.send_document(call.message.chat.id, file_bytes, visible_file_name=f"{topic}.pptx")
    else:
        file_bytes = create_docx(topic)
        bot.send_document(call.message.chat.id, file_bytes, visible_file_name=f"{topic}.docx")
    
    # --- YANGI QO'SHILGAN QISM ---
    # Foydalanuvchiga keyingi qadamni ko'rsatish
    bot.send_message(call.message.chat.id, "✅ Tayyor! Yana boshqa mavzuda yasash uchun shunchaki yangi mavzu nomini yozib yuboring.")
    
    # Senga bildirishnoma yuborish
    bot.send_message(MY_CHAT_ID, f"👤 Foydalanuvchi @{call.from_user.username} {topic} bo'yicha {type_file} yukladi.")

# Botni alohida oqimda ishga tushirish
def run_bot():
    bot.polling(none_stop=True)

if "bot_thread" not in st.session_state:
    thread = threading.Thread(target=run_bot, daemon=True)
    thread.start()
    st.session_state.bot_thread = True

# ==========================================
# 4. STREAMLIT INTERFEYSI (SAYT UCHUN)
# ==========================================
st.set_page_config(page_title="Samir AI", page_icon="🚀")
st.title("🤖 Samir AI Generator")

st.markdown(f"""
    * Bu yerda mavzu yozsangiz, fayllarni yuklab olishingiz mumkin.
    * Telegram botimiz orqali ham foydalanish mumkin: [@wordpower1_bot]
""")

st_topic = st.text_input("Mavzu kiriting:")

if st_topic:
    with st.spinner("Tayyorlanmoqda..."):
        p_bytes = create_pptx(st_topic)
        w_bytes = create_docx(st_topic)
        
        st.success("Tayyor!")
        st.download_button("PowerPoint yuklash", p_bytes, f"{st_topic}.pptx")
        st.download_button("Word yuklash", w_bytes, f"{st_topic}.docx")
        
        bot.send_message(MY_CHAT_ID, f"🌐 Saytdan yangi so'rov: {st_topic}")
